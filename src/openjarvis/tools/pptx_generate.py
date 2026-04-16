"""PPTX generation tool — create PPTX presentation files from text content or a topic."""

from __future__ import annotations

from pathlib import Path
from typing import Any
import json

from openjarvis.core.registry import ToolRegistry
from openjarvis.core.types import ToolResult
from openjarvis.tools._stubs import BaseTool, ToolSpec

@ToolRegistry.register("pptx_generate")
class PPTXGenerateTool(BaseTool):
    """Generate a PPTX file from text content or a topic and serve it for download."""

    tool_id = "pptx_generate"

    @property
    def spec(self) -> ToolSpec:
        return ToolSpec(
            name="pptx_generate",
            description=(
                "Generate a standard PPTX (PowerPoint) slide deck presentation. You can provide either: "
                "(a) a 'topic' and an optional 'title', and the tool will write the content automatically, "
                "or (b) a 'title' and full 'content' text. This tool creates a real PPTX file and opens it in the UI."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "Title of the PPTX presentation.",
                    },
                    "topic": {
                        "type": "string",
                        "description": "The topic/subject for the PPTX. When provided, the tool will auto-generate comprehensive content for the slides. Use this when you don't have the content written yet.",
                    },
                    "content": {
                        "type": "string",
                        "description": "Full text content for the PPTX slides. Use double newlines to separate slides.",
                    },
                },
                "required": [],
            },
            category="media",
            required_capabilities=[],
        )

    def _generate_content_from_topic(self, topic: str, title: str) -> str:
        """Use the LLM engine to generate content about a topic."""
        import re

        # Extract word count constraint from topic string
        wc_match = re.search(
            r"(?:(\d+)\s*words?\s*(?:or\s*less|max(?:imum)?|limit))|(?:(?:under|fewer\s+than|no\s+more\s+than|at\s+most|up\s+to)\s+(\d+)\s*words?)",
            topic, re.IGNORECASE,
        )
        word_limit: int | None = None
        if wc_match:
            word_limit = int(wc_match.group(1) or wc_match.group(2))

        try:
            from openjarvis.core.config import load_config
            from openjarvis.core.types import Message, Role

            config = load_config()
            model = config.intelligence.default_model
            engine_name = config.intelligence.default_engine
            
            if engine_name == "ollama":
                from openjarvis.engine.ollama import OllamaEngine
                engine = OllamaEngine()
            else:
                from openjarvis.engine.cloud import CloudEngine
                engine = CloudEngine()

            constraint_msg = ""
            if word_limit:
                constraint_msg = (
                    f"\n\nCRITICAL CONSTRAINT: Your ENTIRE response must be {word_limit} words or fewer. "
                    "Count carefully. Do NOT exceed this limit."
                )

            messages = [
                Message(role=Role.SYSTEM, content=(
                    "You are a professional presentation creator. Write well-structured content for a slide deck. "
                    "Separate each slide by two blank lines (double newline). For each slide, the first line should be the slide title (starting with #), "
                    "followed by bullet points for the body. KEEP IT CONCISE."
                    f"{constraint_msg}"
                )),
                Message(role=Role.USER, content=(
                    f"Create a presentation titled '{title}' about: {topic}"
                    + (f"\nRemember: maximum {word_limit} words total." if word_limit else "")
                )),
            ]

            gen_max_tokens = min(word_limit * 3, 200) if word_limit else 800

            result = engine.chat(
                messages,
                model=model,
                temperature=config.intelligence.temperature,
                max_tokens=gen_max_tokens,
            )

            # Some models return string directly, others return dict
            if isinstance(result, dict):
                return result.get("content", "")
            elif hasattr(result, "content"):
                return result.content
            return str(result)
        except Exception as e:
            import logging
            logging.getLogger(__name__).error(f"Error generating PPTX content: {e}")
            return f"Topic: {topic}\n\nCould not auto-generate deep content."

    def execute(self, **kwargs: Any) -> ToolResult:
        """Create the PPTX file."""
        title = kwargs.get("title", "")
        topic = kwargs.get("topic", "")
        content = kwargs.get("content", "")

        if not content and topic:
            if not title:
                title = topic[:30] + "..." if len(topic) > 30 else topic
            content = self._generate_content_from_topic(topic, title)

        if not content:
            content = "Presentation\n\nNo content was provided."

        if not title:
            title = "Presentation"

        try:
            from pptx import Presentation
        except ImportError:
            return ToolResult(
                tool_name="pptx_generate",
                content=(
                    "python-pptx package not installed."
                    " Install with: uv pip install python-pptx"
                ),
                success=False,
            )

        try:
            prs = Presentation()
            
            # Sub-split content into slides based on double-newlines or distinct slide markers
            import re
            slides = re.split(r'\n\s*\n', content.strip())
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title_shape.text = title
            subtitle.text = ""
            
            bullet_slide_layout = prs.slide_layouts[1]
            
            # Content slides
            for slide_text in slides:
                if not slide_text.strip():
                    continue
                if title.lower() in slide_text.lower() and len(slide_text) < len(title) + 20:
                    # Skip duplicate title slide
                    continue
                    
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                t_shape = shapes.title
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                
                lines = slide_text.split('\n')
                
                if len(lines) > 0:
                    # Header
                    header = lines[0].strip('# *:_-')
                    t_shape.text = header
                    lines = lines[1:]
                else:
                    t_shape.text = "Slide"
                
                # Add body lines
                first_line = True
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    clean_line = line.lstrip('*-• ')
                    if first_line:
                        tf.text = clean_line
                        first_line = False
                    else:
                        p = tf.add_paragraph()
                        p.text = clean_line
                        # Indent bullet if it was double dashed -- etc.
                        if line.startswith('  -') or line.startswith('  *'):
                            p.level = 1

            # Save to public frontend folder
            frontend_dir = Path("/Users/danny/Desktop/OpenJarvis/frontend/public/generated")
            frontend_dir.mkdir(parents=True, exist_ok=True)

            safe_title = "".join([c if c.isalnum() else "_" for c in title])
            file_path = frontend_dir / f"{safe_title}.pptx"

            prs.save(str(file_path))
            
            public_url = f"/generated/{safe_title}.pptx"

            return ToolResult(
                tool_name="pptx_generate",
                content=f"PPTX generated successfully. The presentation '{title}' has been created and is available for download at {public_url}.\n\n[RENDER_CANVAS:url:{public_url}]\n\nSYSTEM NOTE TO ASSISTANT: You have successfully created the PPTX. Tell the user it's ready for download.",
                success=True,
                metadata={
                    "file_path": str(file_path),
                    "url": public_url
                },
            )
        except Exception as exc:
            return ToolResult(
                tool_name="pptx_generate",
                content=f"PPTX generation error: {exc}",
                success=False,
            )

__all__ = ["PPTXGenerateTool"]
